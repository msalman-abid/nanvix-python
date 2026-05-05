"""Test: python-pptx"""
import sys
sys.stdout.reconfigure(line_buffering=True)
try:
    import lxml  # noqa: F401 — hard dependency of python-pptx
except ImportError:
    print("python-pptx: SKIP (lxml not available)")
    sys.exit(0)
try:
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "NanVix Test"

    assert len(prs.slides) == 1
    assert title.text == "NanVix Test"
    print("python-pptx: PASS")
except Exception as e:
    print(f"python-pptx: FAIL: {e}")
    sys.exit(1)
